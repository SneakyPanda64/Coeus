import collections
import collections.abc
from pptx import Presentation
import os
import openai
from dotenv import load_dotenv

import gpt

load_dotenv()

class Slide:
    def __init__(self, index, paragraphs):
        self.index = index
        self.paragraphs = paragraphs
    def __str__(self):
        return f"Slide({self.index}, {self.paragraphs})"

# def FormatMD(text):
#     openai.api_key = os.getenv("OPENAI_API_KEY")
#     response = openai.ChatCompletion.create(
#         model="gpt-3.5-turbo",
#         max_tokens=2048,
#         temperature=0.2,
#         messages=[
#             {"role": "system", "content": "You are an AI used for teaching business to A level students."},
#             {"role": "user", "content": f"Format the following passage to make it more readable in markdown. And fix the grammar and formatting of the text. Discard and do not include any links.\n\n{text}"},
#         ]
#     )
#     return response.choices[0].message.content
def SimplifyText(text):
    formatted_text = gpt.ChatCompletion([
        {"role": "system", "content": "You are an AI used for teaching business to A level students."},
        {"role": "user",
         "content": f"Loosely simplify the following paragraph:\n\n{text}"},
    ], 2048, 0.2)
    return formatted_text
def FormatText(text):
    formatted_text = gpt.ChatCompletion([
        {"role": "system", "content": "You are an AI used for teaching business to A level students."},
        {"role": "user",
         "content": f"Format the following passage to make it more readable in markdown. And fix the grammar and formatting of the text. Discard and do not include any links. Highlight any important keywords with bold.\n\n{text}"},
    ], 2048, 0.2)
    return formatted_text
def ConvertPP(input, output):
    prs = Presentation(input)
    slides = []
    index = 0
    for slide in prs.slides:
        text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
        if len(text) > 1 and len(" ".join(text)) > 20 and text[0].strip() not in ["Learning Objectives"]:
            slides.append(Slide(index, text))
            index += 1
    with open(output, "w") as f:
        text = ""
        for s in slides:
            body_text = " ".join(s.paragraphs)
            slide_text = f"# {s.paragraphs[0]}\n\n{body_text}\n\n"
            text += slide_text
        formatted_text = FormatText(text)
        f.write(formatted_text)
def ConvertPDF(input, output):
    pass

def Checker():
    out_files = [(os.path.splitext(v)[0]) for v in os.listdir("output")]
    for file in os.listdir("input"):
        if os.path.splitext(file)[0] not in out_files:
            print(f"[Checker] Found file {file} (converting)...")
            file_output = os.path.join("output", os.path.splitext(file)[0] + ".md")
            file_input = os.path.join("input", file)
            match os.path.splitext(file)[1]:
                case "pptx":-
                    ConvertPP(file_input, file_output)
                case "pdf":
                    ConvertPDF(file_input, file_output)
